import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / ".tools" / "python-pptx"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
# La présentation demandée est modifiée directement : le fichier final conserve
# donc le nom utilisé pour la soutenance.
OUTPUT = ROOT / "IOL_Soutenance.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(11, 31, 51)
TEAL = RGBColor(20, 166, 155)
GOLD = RGBColor(245, 178, 69)
BLUE = RGBColor(57, 116, 194)
PURPLE = RGBColor(127, 89, 176)
RED = RGBColor(196, 82, 82)
INK = RGBColor(31, 45, 61)
MUTED = RGBColor(95, 109, 125)
LINE = RGBColor(220, 228, 235)
PALE_TEAL = RGBColor(233, 248, 246)
PALE_BLUE = RGBColor(235, 243, 252)
PALE_GOLD = RGBColor(255, 247, 229)
PALE_RED = RGBColor(253, 240, 240)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(10)
        p.text = "• " + item
    return box


def add_shape(slide, shape_type, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    kind = SHAPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_card(slide, title, body, x, y, w, h, accent=TEAL, fill=WHITE,
             title_size=17, body_size=13):
    shape = add_shape(slide, SHAPE.RECTANGLE, x, y, w, h, fill, LINE, radius=True)
    bar = slide.shapes.add_shape(SHAPE.RECTANGLE, x, y, Inches(0.10), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + Inches(0.22), y + Inches(0.18), w - Inches(0.38), Inches(0.34),
             title_size, accent, True)
    add_text(slide, body, x + Inches(0.22), y + Inches(0.60), w - Inches(0.40), h - Inches(0.72),
             body_size, INK)
    return shape


def add_chevron(slide, x, y, w=Inches(0.34), h=Inches(0.32), color=TEAL):
    chevron = slide.shapes.add_shape(SHAPE.CHEVRON, x, y, w, h)
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = color
    chevron.line.fill.background()
    return chevron


def add_header(slide, section, title, subtitle="", page=1):
    set_background(slide)
    top = slide.shapes.add_shape(SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()
    add_text(slide, section.upper(), Inches(0.60), Inches(0.31), Inches(2.4), Inches(0.25),
             10, TEAL, True)
    add_text(slide, title, Inches(0.60), Inches(0.57), Inches(11.9), Inches(0.54),
             27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.60), Inches(1.16), Inches(11.6), Inches(0.35),
                 13, MUTED)
    footer = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.60), Inches(7.08), Inches(12.13), Inches(0.01))
    footer.fill.solid()
    footer.fill.fore_color.rgb = LINE
    footer.line.fill.background()
    add_text(slide, "IOL — Soutenance de mémoire", Inches(0.60), Inches(7.14), Inches(3.0), Inches(0.18),
             9, MUTED)
    add_text(slide, f"{page:02d}", Inches(12.15), Inches(7.11), Inches(0.55), Inches(0.2),
             10, TEAL, True, PP_ALIGN.RIGHT)


def process_row(slide, labels, y, colors=None, x=Inches(0.70), total_w=Inches(11.95), h=Inches(1.12)):
    count = len(labels)
    gap = Inches(0.12)
    arrow_w = Inches(0.28)
    card_w = (total_w - (count - 1) * (gap + arrow_w)) / count
    current = x
    for i, label in enumerate(labels):
        color = colors[i] if colors else TEAL
        add_shape(slide, SHAPE.RECTANGLE, current, y, card_w, h, WHITE, color, radius=True)
        add_text(slide, label, current + Inches(0.08), y + Inches(0.25), card_w - Inches(0.16), h - Inches(0.35),
                 15, INK, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        current += card_w
        if i < count - 1:
            add_chevron(slide, current + Inches(0.07), y + Inches(0.40), arrow_w, Inches(0.30), colors[i] if colors else TEAL)
            current += arrow_w + gap


def architecture_general(slide):
    add_header(slide, "Architecture", "Architecture générale", "Deux plans séparés : piloter et exécuter.", 7)
    add_card(slide, "Plan de pilotage", "Portail, API, métadonnées, planification, audit et suivi.",
             Inches(0.75), Inches(1.75), Inches(2.55), Inches(1.35), TEAL, PALE_TEAL)
    add_card(slide, "Transport sécurisé", "Source Gateway transporte les données ; Kafka porte les commandes et états ; RustFS prend les gros objets.",
             Inches(3.65), Inches(1.75), Inches(2.85), Inches(1.35), BLUE, PALE_BLUE)
    add_card(slide, "Exécution", "Pipeline Consumer vérifie, matérialise puis lance Hop ou Spark.",
             Inches(6.85), Inches(1.75), Inches(2.40), Inches(1.35), GOLD, PALE_GOLD)
    add_card(slide, "Lakehouse", "Bronze → Silver → Gold dans la destination analytique.",
             Inches(9.60), Inches(1.75), Inches(2.75), Inches(1.35), PURPLE, RGBColor(245, 240, 252))
    for xpos, color in [(Inches(3.35), TEAL), (Inches(6.55), BLUE), (Inches(9.30), GOLD)]:
        add_chevron(slide, xpos, Inches(2.28), Inches(0.25), Inches(0.28), color)
    add_card(slide, "Interopérabilité", "OpenHIM, adaptateurs de normes et format pivot relient les systèmes externes au même cœur de traitement.",
             Inches(1.55), Inches(4.00), Inches(4.15), Inches(1.35), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "Règle essentielle", "Les credentials source restent au Source Gateway. Hop et Spark ne reçoivent qu’un artefact déjà transporté.",
             Inches(6.20), Inches(4.00), Inches(5.55), Inches(1.35), RED, PALE_RED)


def etl_architecture(slide):
    add_header(slide, "Architecture", "Architecture ETL", "Le pipeline est configuré par métadonnées, puis exécuté de manière contrôlée.", 8)
    process_row(slide, ["Configurer\nle workflow", "Transporter\nl’artefact", "Exécuter\nla transformation", "Publier\nles résultats"], Inches(1.85), [TEAL, BLUE, GOLD, PURPLE])
    add_card(slide, "Métadonnées", "Sources, champs, mappings, planning, destination et règles de transformation.",
             Inches(0.90), Inches(3.55), Inches(3.35), Inches(1.45), TEAL, PALE_TEAL)
    add_card(slide, "Contrôles", "Intégrité, droits, suivi d’exécution, journalisation et reprise contrôlée.",
             Inches(4.98), Inches(3.55), Inches(3.35), Inches(1.45), BLUE, PALE_BLUE)
    add_card(slide, "Résultat", "Une exécution produit des statuts et des données exploitables dans les zones du datawarehouse.",
             Inches(9.05), Inches(3.55), Inches(3.35), Inches(1.45), GOLD, PALE_GOLD)
    add_text(slide, "Principe : ajouter une source via les connecteurs et la configuration disponibles ; une nouvelle technologie non supportée exige encore un adaptateur.",
             Inches(0.92), Inches(5.55), Inches(11.3), Inches(0.50), 14, MUTED, False, PP_ALIGN.CENTER)


def business_architecture(slide):
    add_header(slide, "Architecture", "Architecture métier et interopérabilité", "Un vocabulaire pivot réduit la dépendance entre administrations.", 9)
    process_row(slide, ["Système\németteur", "Adaptateur\nde norme", "Format\npivot", "Datawarehouse\net destinataire"], Inches(1.80), [BLUE, PURPLE, TEAL, GOLD])
    add_card(slide, "Entrant", "Réception d’un message, validation du format puis normalisation vers le pivot.",
             Inches(1.10), Inches(3.55), Inches(4.75), Inches(1.35), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "Sortant", "Lecture du résultat Gold, dé-normalisation vers le format attendu et livraison tracée.",
             Inches(7.05), Inches(3.55), Inches(4.15), Inches(1.35), BLUE, PALE_BLUE)
    add_text(slide, "Le pivot porte le sens métier : nom normalisé, type, règle de format et caractère obligatoire.",
             Inches(1.15), Inches(5.50), Inches(11.0), Inches(0.35), 15, NAVY, True, PP_ALIGN.CENTER)


def bigdata_architecture(slide):
    add_header(slide, "Architecture", "Architecture Big Data", "Règle cible de production : privilégier la sécurité lorsque le volume est incertain.", 10)
    add_card(slide, "1. Spark imposé", "Une étape distribuée demandée par le workflow entraîne Spark.",
             Inches(0.80), Inches(1.75), Inches(3.70), Inches(1.20), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "2. Spark préventif", "Estimation absente ou incomplète, ou seuil atteint : Spark.",
             Inches(4.82), Inches(1.75), Inches(3.70), Inches(1.20), GOLD, PALE_GOLD)
    add_card(slide, "3. Local", "Uniquement si toutes les sources sont mesurables et sous les seuils.",
             Inches(8.84), Inches(1.75), Inches(3.70), Inches(1.20), TEAL, PALE_TEAL)
    add_card(slide, "Transport borné", "Lots Kafka et multipart RustFS limitent la mémoire utilisée. Les capacités dépendent aussi du disque, du réseau et des ressources serveurs.",
             Inches(1.00), Inches(3.75), Inches(5.15), Inches(1.45), BLUE, PALE_BLUE)
    add_card(slide, "À qualifier avant GO", "Comparer Hop et Spark sur les mêmes données : schéma, volumes, sommes, dates, valeurs nulles et empreinte canonique.",
             Inches(7.15), Inches(3.75), Inches(5.15), Inches(1.45), RED, PALE_RED)


def warehouse_architecture(slide):
    add_header(slide, "Architecture", "Architecture Datawarehouse", "Modèle cible : données reçues, nettoyées puis consolidées pour la décision.", 11)
    process_row(slide, ["Bronze", "Silver", "Gold"], Inches(1.85), [BLUE, TEAL, GOLD], Inches(1.20), Inches(10.90), Inches(1.25))
    add_text(slide, "Réception brute append-only\n+ identifiant d’exécution\n+ invalidation logique", Inches(1.20), Inches(3.28), Inches(3.10), Inches(0.85), 14, INK, False, PP_ALIGN.CENTER)
    add_text(slide, "Nettoyage et normalisation\npar source ; règles explicites", Inches(5.05), Inches(3.28), Inches(3.10), Inches(0.85), 14, INK, False, PP_ALIGN.CENTER)
    add_text(slide, "Faits et dimensions\npour l’analyse et la restitution", Inches(8.90), Inches(3.28), Inches(3.10), Inches(0.85), 14, INK, False, PP_ALIGN.CENTER)
    add_card(slide, "Grain Gold proposé", "Une ligne représente le total journalier pour un territoire, une nature d’impôt et un type d’usager. Clé : date + territoire + impôt + usager + version de calcul.",
             Inches(1.25), Inches(4.65), Inches(10.80), Inches(1.20), GOLD, PALE_GOLD, 16, 14)


def use_cases(slide):
    add_header(slide, "Conception", "Cas d’utilisation principaux", "Les rôles ont des responsabilités distinctes, contrôlées côté serveur.", 12)
    add_card(slide, "Administrateur", "Gère les utilisateurs, connexions, normes et workflows.",
             Inches(0.80), Inches(2.00), Inches(3.00), Inches(1.30), TEAL, PALE_TEAL)
    add_card(slide, "Plateforme IOL", "Configure\nPlanifie\nExécute\nSupervise\nTrace", 
             Inches(5.15), Inches(1.55), Inches(3.05), Inches(2.20), NAVY, WHITE, 18, 16)
    add_card(slide, "Utilisateur", "Consulte les exécutions, visualise les résultats et teste les requêtes autorisées.",
             Inches(9.55), Inches(2.00), Inches(3.00), Inches(1.30), BLUE, PALE_BLUE)
    add_card(slide, "Systèmes externes", "Envoient ou reçoivent des flux via des canaux d’interopérabilité tracés.",
             Inches(3.20), Inches(4.55), Inches(6.90), Inches(1.15), PURPLE, RGBColor(245, 240, 252))
    add_text(slide, "Administrateur", Inches(3.95), Inches(2.52), Inches(1.0), Inches(0.25), 11, MUTED, False, PP_ALIGN.RIGHT)
    add_chevron(slide, Inches(4.98), Inches(2.52), Inches(0.28), Inches(0.25), TEAL)
    add_chevron(slide, Inches(8.22), Inches(2.52), Inches(0.28), Inches(0.25), BLUE)
    add_text(slide, "Utilisateur", Inches(8.52), Inches(2.52), Inches(0.9), Inches(0.25), 11, MUTED)


def classes(slide):
    add_header(slide, "Conception", "Diagramme de classes simplifié", "Les métadonnées relient la source, les règles, la destination et le suivi.", 13)
    add_card(slide, "WorkflowConfig", "nom\npriorité\nplanification\nGold", Inches(0.75), Inches(1.80), Inches(2.45), Inches(1.55), TEAL, PALE_TEAL)
    add_card(slide, "SourceDefinition", "protocole\nchargement\nwatermark\nchamps", Inches(3.75), Inches(1.80), Inches(2.45), Inches(1.55), BLUE, PALE_BLUE)
    add_card(slide, "DestinationConnection", "type SGBD\nhôte\nbase\ncredential", Inches(6.75), Inches(1.80), Inches(2.45), Inches(1.55), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "ExecutionLog", "statut\nétapes\nmétriques\nwatermarks", Inches(9.75), Inches(1.80), Inches(2.45), Inches(1.55), GOLD, PALE_GOLD)
    add_card(slide, "Standard / StandardTerm", "vocabulaire métier\ntypes\nrègles de format\nréférentiel", Inches(2.25), Inches(4.20), Inches(3.60), Inches(1.50), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "Version de workflow — cible", "instantané des mappings, SQL, contrat source, référentiel et image moteur pour la reproduction historique.",
             Inches(7.25), Inches(4.20), Inches(4.45), Inches(1.50), RED, PALE_RED)
    for xpos in [Inches(3.24), Inches(6.24), Inches(9.24)]:
        add_chevron(slide, xpos, Inches(2.42), Inches(0.28), Inches(0.24), TEAL)


def sequence(slide):
    add_header(slide, "Conception", "Séquence d’exécution", "Un même parcours pour les déclenchements manuels ou planifiés.", 14)
    steps = [
        ("1", "Administrateur / planificateur\ndéclenche le workflow", TEAL),
        ("2", "API\ncrée le journal d’exécution", BLUE),
        ("3", "Source Gateway\ntransporte les données", PURPLE),
        ("4", "Kafka ou RustFS\nporte lots / manifeste", GOLD),
        ("5", "Pipeline Consumer\nvérifie et matérialise", TEAL),
        ("6", "Hop ou Spark\ntraite Bronze → Silver → Gold", BLUE),
        ("7", "Statuts et métriques\nreviennent au portail", PURPLE),
    ]
    x = Inches(0.55)
    y = Inches(2.15)
    for index, (number, label, color) in enumerate(steps):
        add_shape(slide, SHAPE.OVAL, x, y, Inches(0.42), Inches(0.42), color, color)
        add_text(slide, number, x, y + Inches(0.04), Inches(0.42), Inches(0.24), 13, WHITE, True, PP_ALIGN.CENTER)
        add_shape(slide, SHAPE.ROUNDED_RECTANGLE, x + Inches(0.12), y + Inches(0.65), Inches(1.55), Inches(0.98), WHITE, color, radius=True)
        add_text(slide, label, x + Inches(0.20), y + Inches(0.86), Inches(1.38), Inches(0.58), 12, INK, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if index < len(steps) - 1:
            add_chevron(slide, x + Inches(1.78), y + Inches(0.80), Inches(0.24), Inches(0.24), color)
        x += Inches(1.82)
    add_text(slide, "La réussite est publiée seulement après les contrôles d’intégrité ; le point de reprise est mis à jour après succès.",
             Inches(0.90), Inches(5.65), Inches(11.50), Inches(0.40), 14, NAVY, True, PP_ALIGN.CENTER)


def gains(slide):
    add_header(slide, "Valeur apportée", "Gains pour les administrations", "La plateforme vise un meilleur pilotage, pas seulement un meilleur transfert de fichiers.", 15)
    add_card(slide, "Décision plus rapide", "Une vue Gold consolidée alimente tableaux de bord, suivi des recettes et analyses.",
             Inches(0.85), Inches(1.95), Inches(3.70), Inches(2.05), GOLD, PALE_GOLD, 18, 16)
    add_card(slide, "Intégration maîtrisée", "Les flux sont décrits par configuration, standardisés et tracés de bout en bout.",
             Inches(4.82), Inches(1.95), Inches(3.70), Inches(2.05), TEAL, PALE_TEAL, 18, 16)
    add_card(slide, "Gouvernance renforcée", "Droits, journaux, états d’exécution et contrôles rendent les opérations vérifiables.",
             Inches(8.79), Inches(1.95), Inches(3.70), Inches(2.05), BLUE, PALE_BLUE, 18, 16)
    add_text(slide, "Résultat attendu : moins de consolidation manuelle, une information plus fiable et des échanges plus simples entre systèmes.",
             Inches(1.15), Inches(5.15), Inches(11.0), Inches(0.48), 16, NAVY, True, PP_ALIGN.CENTER)


def limits(slide):
    add_header(slide, "Limites", "Limites et précautions", "Les garanties doivent être démontrées avant d’être annoncées comme acquises.", 16)
    add_card(slide, "Qualification Big Data", "Mesurer mémoire, débit, p95/p99, pannes et résultats Hop/Spark sur plusieurs volumes.",
             Inches(0.80), Inches(1.78), Inches(3.70), Inches(1.65), RED, PALE_RED)
    add_card(slide, "Historique exact", "Mettre en place les versions immuables des règles, contrats et référentiels.",
             Inches(4.82), Inches(1.78), Inches(3.70), Inches(1.65), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "Production", "Finaliser secrets, sauvegardes, restauration, supervision et procédures d’incident.",
             Inches(8.84), Inches(1.78), Inches(3.70), Inches(1.65), BLUE, PALE_BLUE)
    add_card(slide, "Assistant IA", "C’est une aide à la proposition SQL : le résultat reste contrôlé, validé et soumis à l’utilisateur.",
             Inches(2.85), Inches(4.45), Inches(7.65), Inches(1.30), GOLD, PALE_GOLD, 16, 15)


def cover(slide):
    set_background(slide, NAVY)
    stripe = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.70), Inches(0.80), Inches(0.10), Inches(4.90))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()
    add_text(slide, "SOUTENANCE DE MÉMOIRE", Inches(1.15), Inches(0.95), Inches(5.0), Inches(0.25), 12, TEAL, True)
    add_text(slide, "Conception et mise en œuvre d’une plateforme ETL interopérable", Inches(1.15), Inches(1.45), Inches(10.6), Inches(1.25), 31, WHITE, True)
    add_text(slide, "Pilotée par métadonnées pour la construction d’un datawarehouse des services publics", Inches(1.15), Inches(2.90), Inches(9.6), Inches(0.75), 19, RGBColor(215, 228, 240))
    add_text(slide, "Présenté par Ngono Atenga Arnaud Andy", Inches(1.15), Inches(5.70), Inches(5.5), Inches(0.32), 15, WHITE, True)
    add_text(slide, "Année académique 2025–2026", Inches(1.15), Inches(6.10), Inches(4.5), Inches(0.25), 12, RGBColor(185, 204, 220))
    add_shape(slide, SHAPE.ROUNDED_RECTANGLE, Inches(9.90), Inches(5.35), Inches(2.30), Inches(0.70), PALE_GOLD, GOLD, radius=True)
    add_text(slide, "TresorPay", Inches(10.05), Inches(5.56), Inches(2.0), Inches(0.22), 18, NAVY, True, PP_ALIGN.CENTER)


def plan(slide):
    add_header(slide, "Plan", "Un parcours simple", "Du besoin des administrations à une solution qualifiée pour la production.", 2)
    entries = [
        ("01", "Contexte", TEAL), ("02", "Problème", RED), ("03", "Solution", GOLD),
        ("04", "Architectures", BLUE), ("05", "Conception", PURPLE), ("06", "Gains et limites", TEAL),
    ]
    for i, (number, label, color) in enumerate(entries):
        col, row = i % 3, i // 3
        x, y = Inches(0.90 + col * 4.15), Inches(1.85 + row * 2.15)
        add_shape(slide, SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.55), Inches(1.48), WHITE, color, radius=True)
        add_text(slide, number, x + Inches(0.25), y + Inches(0.26), Inches(0.70), Inches(0.35), 23, color, True)
        add_text(slide, label, x + Inches(1.05), y + Inches(0.38), Inches(2.15), Inches(0.30), 19, NAVY, True)


def context(slide):
    add_header(slide, "Contexte", "Des données publiques nombreuses, mais dispersées", "Le besoin est né du terrain : mieux exploiter les données administratives.", 3)
    add_card(slide, "Sources hétérogènes", "Bases SQL et NoSQL, fichiers, API et applications métier coexistent.",
             Inches(0.85), Inches(2.00), Inches(3.70), Inches(2.15), BLUE, PALE_BLUE, 19, 16)
    add_card(slide, "Sens métier différent", "Une même information peut porter un nom, un format ou une règle différente selon l’administration.",
             Inches(4.82), Inches(2.00), Inches(3.70), Inches(2.15), PURPLE, RGBColor(245, 240, 252), 19, 16)
    add_card(slide, "Décision difficile", "Sans consolidation fiable, le reporting reste lent, manuel et incomplet.",
             Inches(8.79), Inches(2.00), Inches(3.70), Inches(2.15), GOLD, PALE_GOLD, 19, 16)
    add_text(slide, "Enjeu : transformer ces données isolées en information fiable pour le pilotage public.",
             Inches(1.00), Inches(5.25), Inches(11.3), Inches(0.40), 17, NAVY, True, PP_ALIGN.CENTER)


def problem(slide):
    add_header(slide, "Problème", "Quatre difficultés à résoudre ensemble", "Le défi ne porte pas seulement sur l’extraction : il concerne aussi le sens, la gouvernance et l’échelle.", 4)
    entries = [
        ("Technique", "Chaque système expose son propre protocole et format.", BLUE),
        ("Sémantique", "Les données comparables n’utilisent pas le même vocabulaire.", PURPLE),
        ("Intégration", "Le point-à-point augmente fortement le coût de chaque nouvelle connexion.", RED),
        ("Gouvernance", "Les données sensibles exigent traçabilité, droits et contrôle des opérations.", GOLD),
    ]
    for i, (title, body, color) in enumerate(entries):
        col, row = i % 2, i // 2
        add_card(slide, title, body, Inches(1.00 + col * 5.85), Inches(1.85 + row * 2.10), Inches(5.00), Inches(1.45), color,
                 PALE_BLUE if color == BLUE else PALE_GOLD if color == GOLD else RGBColor(245, 240, 252) if color == PURPLE else PALE_RED,
                 18, 15)


def question(slide):
    add_header(slide, "Problème", "Question et objectif", "Une plateforme unique doit faciliter l’intégration sans réduire les exigences de contrôle.", 5)
    add_shape(slide, SHAPE.ROUNDED_RECTANGLE, Inches(1.00), Inches(1.75), Inches(11.30), Inches(1.65), PALE_TEAL, TEAL, radius=True)
    add_text(slide, "Comment intégrer, normaliser et consolider des flux publics hétérogènes sans reconstruire un pipeline pour chaque source, tout en assurant sécurité, traçabilité et capacité d’évolution ?",
             Inches(1.35), Inches(2.12), Inches(10.60), Inches(0.85), 21, NAVY, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    process_row(slide, ["Intégrer", "Normaliser", "Consolider", "Gouverner"], Inches(4.35), [BLUE, PURPLE, GOLD, TEAL], Inches(1.20), Inches(10.90), Inches(0.90))


def solution(slide):
    add_header(slide, "Solution", "La réponse apportée par IOL", "Une chaîne d’intégration guidée par les métadonnées, le pivot métier et le raffinage progressif.", 6)
    process_row(slide, ["Configurer", "Transporter", "Raffiner", "Consolider", "Restituer"], Inches(1.85), [TEAL, BLUE, PURPLE, GOLD, TEAL], Inches(0.65), Inches(12.05), Inches(1.15))
    add_card(slide, "Pilotage par métadonnées", "Les paramètres du flux sont stockés et contrôlés : source, champs, mappings, règles et planning.",
             Inches(0.85), Inches(3.70), Inches(3.55), Inches(1.45), TEAL, PALE_TEAL)
    add_card(slide, "Format pivot", "Un vocabulaire commun relie les systèmes qui ne partagent ni technologie ni modèle métier.",
             Inches(4.90), Inches(3.70), Inches(3.55), Inches(1.45), PURPLE, RGBColor(245, 240, 252))
    add_card(slide, "Datawarehouse", "Les données progressent de Bronze vers Silver puis Gold, en fonction de leur niveau de qualité et d’usage.",
             Inches(8.95), Inches(3.70), Inches(3.55), Inches(1.45), GOLD, PALE_GOLD)


def conclusion(slide):
    add_header(slide, "Conclusion", "Ce que la plateforme apporte", "Une base d’intégration et de décision qui reste à qualifier complètement en conditions réelles.", 17)
    add_text(slide, "IOL fait passer les administrations de données isolées à des flux configurés, tracés et consolidés.",
             Inches(1.10), Inches(1.75), Inches(11.1), Inches(0.55), 24, NAVY, True, PP_ALIGN.CENTER)
    process_row(slide, ["Sources\nhétérogènes", "Pivot\nmétier", "Datawarehouse\nexploitable"], Inches(3.15), [BLUE, TEAL, GOLD], Inches(1.65), Inches(9.90), Inches(1.15))
    add_text(slide, "Prochaine étape : finaliser les corrections de production, exécuter les qualifications de charge et de reprise, puis décider du GO/NO-GO.",
             Inches(1.10), Inches(5.35), Inches(11.1), Inches(0.45), 16, MUTED, False, PP_ALIGN.CENTER)


def thanks(slide):
    set_background(slide, NAVY)
    add_text(slide, "Merci pour votre attention", Inches(1.10), Inches(2.35), Inches(11.1), Inches(0.70), 34, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "Questions et échanges", Inches(1.10), Inches(3.35), Inches(11.1), Inches(0.35), 19, TEAL, True, PP_ALIGN.CENTER)
    add_shape(slide, SHAPE.ROUNDED_RECTANGLE, Inches(5.38), Inches(5.35), Inches(2.55), Inches(0.66), PALE_GOLD, GOLD, radius=True)
    add_text(slide, "IOL / TresorPay", Inches(5.50), Inches(5.55), Inches(2.30), Inches(0.20), 15, NAVY, True, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    builders = [
        cover, plan, context, problem, question, solution,
        architecture_general, etl_architecture, business_architecture,
        bigdata_architecture, warehouse_architecture, use_cases, classes,
        sequence, gains, limits, conclusion, thanks,
    ]
    for builder in builders:
        builder(prs.slides.add_slide(blank))
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
